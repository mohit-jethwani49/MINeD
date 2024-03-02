from flask import Flask, render_template, request, jsonify, session
import os
import fitz  # PyMuPDF
import string

from langchain_google_genai import ChatGoogleGenerativeAI
import os
import google.generativeai as genai
from pptx import Presentation

import aspose.pdf as ap
from docx import Document
# from flask import Flask, render_template



def gemini(text,query):
    os.environ['GOOGLE_API_KEY'] = "AIzaSyBRpB3U0FbKL_m10g-q1e1LBuPVZBkG7dc"
    genai.configure(api_key = os.environ['GOOGLE_API_KEY'])
    llm = ChatGoogleGenerativeAI(model="gemini-pro")
    def Generate_Response(text, query):
        prompt = "'''\n" + text +"\n'''" + "Task: \n\nAnswer the query strictly based on the above text:\n" + query
        response = llm.invoke(prompt)
        return (response.content)
    return Generate_Response(text,query)

def pdf(filename):
    filename='uploads\\'+filename
    def extract_ascii(text):
        ascii_chars = set(string.ascii_letters + string.punctuation + string.digits + string.whitespace)
        ascii_text = ''.join(char for char in text if char in ascii_chars)
        return ascii_text

    def get_text(pdf_file_path):
        # Open the PDF file
        pdf_document = fitz.open(pdf_file_path)

        text_content = ''

        for page_number in range(pdf_document.page_count):
            # Get a page
            page = pdf_document[page_number]

            # Extract text from the page
            text_content += page.get_text()
        text= extract_ascii(text_content)
        text= text.replace("| link","")
        text=text.replace("| Link","")
        text=text.replace("#","")
        text=text.replace("|","")
        # print(text)
        # Close the PDF file
        
        pdf_document.close()

        # Write the extracted text to a text file
        with open("output_a.txt", 'w', encoding='utf-8') as text_file:
            text_file.write(text)
        return text

    # Example usage
    pdf_file_path = filename  
    print(filename)# Replace with your PDF file path
    # Replace with the desired output text file path
    return get_text(pdf_file_path)
    
        
def pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)


def aspos(file):
    # Create TeXLoadOptions class object
    options = ap.TeXLoadOptions()

    # Create a Document class object
    document = ap.Document(file , options)

    # Convert Latex to PDF
    document.save("tex-to-pdf.pdf")
    return pdf("tex-to-pdf.pdf")

def ddd(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

    
    






app = Flask(__name__)

@app.route('/')
def index():
    # return render_template('index.html', my_string="")
    if 'my_string' not in session:
        session['my_string'] = ""

    return render_template('index.html')

@app.route('/process_form', methods=['POST'])
def process_form():
    # Get the query from the request form data
    query = request.form.get('query')

    # Get the file from the request files
    file = request.files['document']
    filename = file.filename

    # Check file extension
    allowed_extensions = {'docx', 'tex', 'pdf', 'pptx'}
    file_extension = filename.rsplit('.', 1)[-1].lower()

    if file_extension not in allowed_extensions:
        return jsonify({'error': 'Unsupported file format'})

    # Save the file to a desired location
    print(filename)
    file.save(os.path.join('uploads', filename))
    textt=""
    if (file_extension=='pdf'):
        textt=pdf(filename)
    elif (file_extension=='pptx'):
        textt=pptx("uploads\\"+filename)
    elif (file_extension=='tex'):
        textt=aspos("uploads\\"+filename)
    elif (file_extension=='docx'):
        textt=ddd("uploads\\"+filename)
        
    else:
        return jsonify({'error': 'Unsupported file format'})
    
    res= gemini(textt,query)
    print(res)
    # return render_template('index.html', my_string=my_string)  
    session['my_string'] = res  # Set your actual string here

    return render_template('index.html')
    
    

    # Process the query and file as needed
    # For demonstration, let's just print them
    print("Query:", query)
    print("File saved as:", filename)

    # You can also send a response back to the client if needed
    return jsonify({'message': 'Submitted successfully'})

if __name__ == '__main__':
    app.secret_key = os.urandom(24)  
    app.run(debug=True)
